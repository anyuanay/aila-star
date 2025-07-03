from fastapi import FastAPI, HTTPException
from pydantic import BaseModel
from fastapi.middleware.cors import CORSMiddleware
from dotenv import load_dotenv
import os
import fitz  # PyMuPDF
from supabase import create_client, Client
from llama_index.llms.gemini import Gemini
from pptx import Presentation

# Load environment variables
load_dotenv()
SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_KEY = os.getenv("SUPABASE_SERVICE_ROLE_KEY")
google_api_key = os.getenv("GOOGLE_API_KEY")
if google_api_key:
    os.environ["GOOGLE_API_KEY"] = google_api_key

# Initialize LLM
llm = Gemini(model="models/gemini-2.5-pro")

class ProcessLectureRequest(BaseModel):
    filename: str
    course_id: str

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

supabase: Client = create_client(SUPABASE_URL, SUPABASE_KEY)

def download_from_supabase(bucket: str, file_path: str, save_as: str):
    response = supabase.storage.from_(bucket).download(file_path)
    if not response:
        raise HTTPException(status_code=404, detail="File not found in Supabase Storage")
    with open(save_as, "wb") as f:
        f.write(response)

def extract_text_segments(pdf_path):
    doc = fitz.open(pdf_path)
    segments = []
    for page in doc:
        text = page.get_text(sort=True)
        if text.strip():
            segments.append(text)
    return segments

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    segments = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        if slide_text:
            segments.append("\n".join(slide_text))
    return segments

@app.post("/process-lecture/")
async def process_lecture(request: ProcessLectureRequest):
    filename = request.filename
    course_id = request.course_id
    bucket = "lecture-materials"
    supabase_path = f"{course_id}/{filename}"
    local_path = f"uploads/{course_id}_{filename}"
    os.makedirs("uploads", exist_ok=True)
    try:
        download_from_supabase(bucket, supabase_path, local_path)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Download failed: {str(e)}")

    # Choose extractor based on file type
    if filename.lower().endswith('.pdf'):
        segments = extract_text_segments(local_path)
    elif filename.lower().endswith('.pptx'):
        segments = extract_text_from_pptx(local_path)
    else:
        raise HTTPException(status_code=400, detail="Unsupported file type")

    # Summarize each segment with Gemini LLM
    summaries = []
    for segment in segments:
        if segment.strip():
            resp = llm.complete(f"Summarize this lecture segment:\n{segment}")
            summaries.append(str(resp))
        else:
            summaries.append("")

    return {
        "status": "processed",
        "filename": filename,
        "num_segments": len(segments),
        "summaries": summaries,
        "segments": segments,  # Return all segments for full UI display
    }

@app.get("/")
async def root():
    return {"message": "AILA Backend is running"}
